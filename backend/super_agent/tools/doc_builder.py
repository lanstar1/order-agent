"""
문서 생성 도구 — docx, pptx, xlsx, pdf 생성
"""
import os
import json
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
from datetime import datetime

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(__file__).parent.parent.parent.parent / "data" / "sa_outputs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def build_document(
    artifact_type: str,
    content: Dict[str, Any],
    filename_prefix: str = "report",
) -> Dict[str, str]:
    """
    문서 생성 메인 함수
    Returns: {file_path, file_name, format}
    """
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    if artifact_type == "report":
        return _build_docx(content, f"{filename_prefix}_{timestamp}")
    elif artifact_type == "slides":
        return _build_pptx(content, f"{filename_prefix}_{timestamp}")
    elif artifact_type == "sheet":
        return _build_xlsx(content, f"{filename_prefix}_{timestamp}")
    elif artifact_type == "brief":
        return _build_markdown(content, f"{filename_prefix}_{timestamp}")
    else:
        return _build_markdown(content, f"{filename_prefix}_{timestamp}")


def _build_docx(content: Dict[str, Any], filename: str) -> Dict[str, str]:
    """Word 문서 생성"""
    try:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # 스타일 설정
        style = doc.styles['Normal']
        font = style.font
        font.name = 'Malgun Gothic'
        font.size = Pt(10)

        # 제목
        title = content.get("title", "분석 보고서")
        heading = doc.add_heading(title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 생성 일시
        date_p = doc.add_paragraph()
        date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = date_p.add_run(f"생성일: {datetime.now().strftime('%Y년 %m월 %d일')}")
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(128, 128, 128)

        doc.add_paragraph()  # 빈 줄

        # 요약
        if content.get("executive_summary"):
            doc.add_heading("Executive Summary", level=1)
            doc.add_paragraph(content["executive_summary"])

        # 섹션들
        for section in content.get("sections", []):
            doc.add_heading(section.get("heading", ""), level=2)
            body = section.get("body", "")
            if isinstance(body, list):
                for item in body:
                    doc.add_paragraph(str(item), style='List Bullet')
            else:
                doc.add_paragraph(str(body))

        # 테이블
        for table_data in content.get("tables", []):
            table_title = table_data.get("title", "")
            if table_title:
                doc.add_heading(table_title, level=3)

            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])
            if headers and rows:
                table = doc.add_table(rows=1 + len(rows), cols=len(headers))
                table.style = 'Table Grid'

                # 헤더
                for i, h in enumerate(headers):
                    cell = table.rows[0].cells[i]
                    cell.text = str(h)
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.bold = True

                # 데이터
                for r_idx, row in enumerate(rows):
                    for c_idx, val in enumerate(row[:len(headers)]):
                        table.rows[r_idx + 1].cells[c_idx].text = str(val)

        # 액션 아이템
        if content.get("action_items"):
            doc.add_heading("액션 아이템", level=1)
            for item in content["action_items"]:
                doc.add_paragraph(str(item), style='List Number')

        # 출처
        if content.get("references"):
            doc.add_heading("참고 자료", level=1)
            for ref in content["references"]:
                doc.add_paragraph(str(ref), style='List Bullet')

        file_name = f"{filename}.docx"
        file_path = str(OUTPUT_DIR / file_name)
        doc.save(file_path)

        return {"file_path": file_path, "file_name": file_name, "format": "docx"}

    except ImportError:
        logger.warning("[DocBuilder] python-docx 미설치, 마크다운으로 대체")
        return _build_markdown(content, filename)


def _build_pptx(content: Dict[str, Any], filename: str) -> Dict[str, str]:
    """PowerPoint 생성"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 타이틀 슬라이드
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = content.get("title", "분석 보고서")
        subtitle.text = f"생성일: {datetime.now().strftime('%Y년 %m월 %d일')}"

        # 요약 슬라이드
        if content.get("executive_summary"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Executive Summary"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = content["executive_summary"]

        # 섹션 슬라이드
        for section in content.get("sections", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section.get("heading", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            body_content = section.get("body", "")
            if isinstance(body_content, list):
                tf.text = body_content[0] if body_content else ""
                for item in body_content[1:]:
                    p = tf.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                tf.text = str(body_content)

            # 슬라이드 노트
            if section.get("notes"):
                slide.notes_slide.notes_text_frame.text = section["notes"]

        # 액션 아이템 슬라이드
        if content.get("action_items"):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "액션 아이템"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = content["action_items"][0] if content["action_items"] else ""
            for item in content["action_items"][1:]:
                p = tf.add_paragraph()
                p.text = str(item)

        file_name = f"{filename}.pptx"
        file_path = str(OUTPUT_DIR / file_name)
        prs.save(file_path)

        return {"file_path": file_path, "file_name": file_name, "format": "pptx"}

    except ImportError:
        logger.warning("[DocBuilder] python-pptx 미설치, 마크다운으로 대체")
        return _build_markdown(content, filename)


def _build_xlsx(content: Dict[str, Any], filename: str) -> Dict[str, str]:
    """Excel 생성"""
    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side

        wb = openpyxl.Workbook()

        # 요약 시트
        ws = wb.active
        ws.title = "요약"

        title = content.get("title", "분석 보고서")
        ws['A1'] = title
        ws['A1'].font = Font(size=16, bold=True)
        ws['A2'] = f"생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        ws['A2'].font = Font(size=10, color="808080")

        row_idx = 4
        if content.get("executive_summary"):
            ws.cell(row=row_idx, column=1, value="Executive Summary").font = Font(size=12, bold=True)
            row_idx += 1
            ws.cell(row=row_idx, column=1, value=content["executive_summary"])
            row_idx += 2

        # 테이블 데이터 시트
        for t_idx, table_data in enumerate(content.get("tables", [])):
            sheet_name = table_data.get("title", f"데이터{t_idx + 1}")[:31]
            ws_data = wb.create_sheet(title=sheet_name)

            headers = table_data.get("headers", [])
            rows = table_data.get("rows", [])

            # 헤더 스타일
            header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True, size=10)

            for c_idx, h in enumerate(headers, 1):
                cell = ws_data.cell(row=1, column=c_idx, value=str(h))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')

            # 데이터
            for r_idx, row in enumerate(rows, 2):
                for c_idx, val in enumerate(row[:len(headers)], 1):
                    cell = ws_data.cell(row=r_idx, column=c_idx)
                    # 숫자 변환 시도
                    try:
                        cell.value = float(str(val).replace(',', ''))
                        cell.number_format = '#,##0'
                    except (ValueError, TypeError):
                        cell.value = str(val)

            # 자동 열 너비
            for col in ws_data.columns:
                max_len = max(len(str(cell.value or "")) for cell in col)
                ws_data.column_dimensions[col[0].column_letter].width = min(max_len + 2, 40)

        file_name = f"{filename}.xlsx"
        file_path = str(OUTPUT_DIR / file_name)
        wb.save(file_path)

        return {"file_path": file_path, "file_name": file_name, "format": "xlsx"}

    except ImportError:
        logger.warning("[DocBuilder] openpyxl 미설치, 마크다운으로 대체")
        return _build_markdown(content, filename)


def _build_markdown(content: Dict[str, Any], filename: str) -> Dict[str, str]:
    """마크다운 문서 생성 (항상 가능한 폴백)"""
    lines = []
    title = content.get("title", "분석 보고서")
    lines.append(f"# {title}")
    lines.append(f"*생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}*")
    lines.append("")

    if content.get("executive_summary"):
        lines.append("## Executive Summary")
        lines.append(content["executive_summary"])
        lines.append("")

    for section in content.get("sections", []):
        lines.append(f"## {section.get('heading', '')}")
        body = section.get("body", "")
        if isinstance(body, list):
            for item in body:
                lines.append(f"- {item}")
        else:
            lines.append(str(body))
        lines.append("")

    for table_data in content.get("tables", []):
        if table_data.get("title"):
            lines.append(f"### {table_data['title']}")
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if headers:
            lines.append("| " + " | ".join(str(h) for h in headers) + " |")
            lines.append("| " + " | ".join("---" for _ in headers) + " |")
            for row in rows:
                lines.append("| " + " | ".join(str(v) for v in row[:len(headers)]) + " |")
            lines.append("")

    if content.get("action_items"):
        lines.append("## 액션 아이템")
        for i, item in enumerate(content["action_items"], 1):
            lines.append(f"{i}. {item}")
        lines.append("")

    text = '\n'.join(lines)
    file_name = f"{filename}.md"
    file_path = str(OUTPUT_DIR / file_name)
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(text)

    return {"file_path": file_path, "file_name": file_name, "format": "markdown"}
